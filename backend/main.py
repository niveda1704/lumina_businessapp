from fastapi import FastAPI, HTTPException, Depends
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel
from typing import List, Optional
import random
import os
import json
from datetime import datetime
from sqlalchemy import create_engine, Column, Integer, String, Float, Text, JSON, ForeignKey
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker, Session, relationship
from passlib.context import CryptContext
from jose import jwt, JWTError
from datetime import datetime, timedelta
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Load Environment Variables
load_dotenv()

app = FastAPI(title="LUMINA Operational Intelligence")

# Database Setup
SQLALCHEMY_DATABASE_URL = "sqlite:///./ibld.db"
engine = create_engine(SQLALCHEMY_DATABASE_URL, connect_args={"check_same_thread": False})
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()

def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

# --- AUTH CONFIG ---
SECRET_KEY = os.getenv("SECRET_KEY", "fallback_insecure_dev_key")
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 60 * 24

pwd_context = CryptContext(schemes=["pbkdf2_sha256"], deprecated="auto")
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="token")

class UserDB(Base):
    __tablename__ = "users"
    id = Column(Integer, primary_key=True, index=True)
    username = Column(String, unique=True, index=True)
    hashed_password = Column(String)
    role = Column(String, default="consultant")

class WorkflowDB(Base):
    __tablename__ = "workflows"
    id = Column(Integer, primary_key=True, index=True)
    name = Column(String, index=True)
    description = Column(String)
    created_at = Column(String)
    input_data = Column(JSON)
    result_data = Column(JSON)
    owner_id = Column(Integer, ForeignKey("users.id"), nullable=True)

Base.metadata.create_all(bind=engine)

# --- AUTH HELPERS ---
def verify_password(plain, hashed): return pwd_context.verify(plain, hashed)
def get_password_hash(password): return pwd_context.hash(password)
def create_access_token(data: dict):
    to_encode = data.copy()
    expire = datetime.utcnow() + timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    to_encode.update({"exp": expire})
    return jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)

def get_current_user(token: str = Depends(oauth2_scheme), db: Session = Depends(get_db)):
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        username: str = payload.get("sub")
        if username is None: raise HTTPException(status_code=401)
    except JWTError: raise HTTPException(status_code=401)
    user = db.query(UserDB).filter(UserDB.username == username).first()
    if user is None: raise HTTPException(status_code=401)
    return user

# CORS Setup
origins = ["*"]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class WorkflowInput(BaseModel):
    name: str
    description: str
    people_involved: int
    approvals_per_task: int
    tools_used: List[str]
    avg_delays_hours: float
    rejection_rate: float = 0.0
    monthly_volume: int = 1 # Default to 1 run per month if unspecified
    avg_annual_salary: float = 2500000.0 # Used for hourly rate calc only
    total_project_budget: float = 5000000.0 # The total pot of money

class LossAnalysis(BaseModel):
    id: Optional[int] = None
    weekly_time_loss_hours: float
    estimated_financial_loss: float
    severity: str
    clarity_score: int
    industry_benchmark_score: int
    decision_delay_index: float
    rework_loss_hours: float
    waste_ratio: float
    total_investment: float
    confidence_interval: dict
    invisible_loss_points: List[dict]
    recommendations: List[dict]

@app.get("/workflows")
def get_workflows(db: Session = Depends(get_db)):
    return db.query(WorkflowDB).all()

@app.post("/analyze", response_model=LossAnalysis)
async def analyze_workflow(data: WorkflowInput, db: Session = Depends(get_db), current_user: UserDB = Depends(get_current_user)):
    # Core Logic for Estimation & Benchmarking
    
    # 1. Benchmarking (Mock Data per Industry Standard)
    industry_benchmark_score = random.randint(65, 85)
    
    # ... existing calculations ...
    # THESE ARE PER-RUN CALCULATIONS
    tool_friction = max(0, len(data.tools_used) - 2) * 0.5 * data.people_involved
    approval_friction = data.approvals_per_task * data.people_involved * 1.5 
    base_delay_impact = data.avg_delays_hours * data.people_involved

    # NEW: Rework Logic
    rework_impact = (base_delay_impact + approval_friction) * (data.rejection_rate / 100) * 1.5
    
    # Total Time Loss PER RUN
    time_loss_per_run = tool_friction + approval_friction + base_delay_impact + rework_impact
    
    # 2. Financial Loss & Investment Analysis
    # Hourly Rate = Annual Salary / (50 weeks * 40 hours) -> ~2000 hours
    hourly_rate = data.avg_annual_salary / 2000
    
    # Financial Loss PER RUN
    financial_loss_per_run = time_loss_per_run * hourly_rate
    
    # SCALING BY VOLUME
    # We assume 'monthly_volume' is how often this process runs per month.
    # Weekly loss = (Per Run * Monthly Volume) / 4 (approx)
    # Annual loss = (Per Run * Monthly Volume) * 12
    
    annual_financial_loss = financial_loss_per_run * data.monthly_volume * 12
    financial_loss_weekly = annual_financial_loss / 50 # Standardized weekly view for dashboard consistency
    
    # Update returned metric name for clarity (variable re-use)
    # weekly_time_loss is actually display variable
    weekly_time_loss = (time_loss_per_run * data.monthly_volume * 12) / 50

    # Total Investment Logic (Fixed per user request)
    # The user enters the specific "Total Project Budget". We trust this number as the Investment Base.
    # If they enter 0, we fallback to estimating it via personnel cost.
    estimated_personnel_cost = data.avg_annual_salary * data.people_involved
    
    total_annual_investment = data.total_project_budget
    if total_annual_investment <= 0:
        total_annual_investment = estimated_personnel_cost # Fallback

    # Annualized Loss (Already calculated above)
    # annual_financial_loss = financial_loss_weekly * 50  <-- We now derive this FROM volume
    
    # Waste Ratio
    waste_ratio = 0
    if total_annual_investment > 0:
        waste_ratio = round((annual_financial_loss / total_annual_investment) * 100, 1)

    # 3. Severity (Dynamic based on % of Investment Wasted)
    if financial_loss_weekly > 500000:
        severity = "High"
    elif financial_loss_weekly > 100000:
        severity = "Medium"
    else:
        severity = "Low"
        
    # 4. Clarity Score (0-100)
    # Deduct points for high people, high tools, high approvals
    clarity = 100
    clarity -= (data.people_involved * 2)
    clarity -= (len(data.tools_used) * 5)
    clarity -= (data.approvals_per_task * 10)
    clarity = max(10, min(100, clarity))
    
    # 5. Invisible Loss Points & AI Reasoning (Heuristic Simulation of Consultant Logic)
    loss_points = []
    
    # Analyze Tool Fragmentation
    if len(data.tools_used) > 3:
        loss_points.append({
            "title": "Cognitive Context Switching",
            "reason": f"Workflow spans {len(data.tools_used)} distinct platforms ({', '.join(data.tools_used[:3])}...), forcing frequent mental re-calibration.",
            "impact": "Reduces deep-work capacity by ~20% per additional tool; increases error rate."
        })
    elif len(data.tools_used) == 0:
        loss_points.append({
            "title": "Digital Opacity",
            "reason": "Lack of defined digital toolset implies reliance on unstructured verbal/manual transmission.",
            "impact": "High risk of tribal knowledge loss and zero auditability."
        })

    # Analyze Approval Friction
    if data.approvals_per_task > 2:
        loss_points.append({
            "title": "Decision Latency Accumulation",
            "reason": f"Requires {data.approvals_per_task} distinct sign-offs. If sequential, this creates exponential delay chains.",
            "impact": "cycle time extends 4x beyond actual working hours due to asynchronous availability."
        })
    
    # Analyze Team Scale vs Complexity
    if data.people_involved > 6:
        loss_points.append({
            "title": "Communication Overhead Entropy",
            "reason": f"{data.people_involved} active participants creates potentially {data.people_involved * (data.people_involved - 1) // 2} communication pathways.",
            "impact": "Information asymmetry increases; projected 15% of time lost clarifyng requirements."
        })
    
    # Analyze Delays
    if data.avg_delays_hours > 4:
         loss_points.append({
            "title": "Idle Capital Detection",
            "reason": "Significant 'wait states' detected between process steps.",
            "impact": "Resources are allocated but blocked; effectively paying for downtime."
        })

    # Feature: Decision Delay Index (DDI)
    # Ratio of "Wait Time" to "Total Cycle Time". Higher is worse.
    # Estimate active work time per person per task ~ 2 hours (heuristic)
    active_work_time = data.people_involved * 2 
    total_wait_time = (data.avg_delays_hours * data.people_involved) + (data.approvals_per_task * 24)
    total_cycle_time = active_work_time + total_wait_time
    
    decision_delay_index = round((total_wait_time / total_cycle_time) * 10, 1) if total_cycle_time > 0 else 0
    decision_delay_index = min(10, decision_delay_index)

    # Feature: Confidence Bands
    # Standard deviation heuristic +/- 15% for ambiguity
    loss_lower_bound = financial_loss_weekly * 0.85
    loss_upper_bound = financial_loss_weekly * 1.15

    # 5. Invisible Loss Points (Enhanced with Root Causes & Blindness)
    loss_points = []
    
    if len(data.tools_used) > 3:
        loss_points.append({
            "title": "Cognitive Context Switching",
            "reason": f"Workflow spans {len(data.tools_used)} distinct platforms, forcing mental re-calibration.",
            "root_cause": "Fragmented IT procurement allowed departments to adopt isolated tools without integration strategy.",
            "blindness_reason": "License costs are visible in budgets, but 'attention residue' time loss does not appear on P&L.",
            "impact": "Reduces deep-work capacity by ~20%."
        })
    elif len(data.tools_used) == 0:
         loss_points.append({
            "title": "Digital Opacity",
            "reason": "Process relies entirely on verbal/manual transmission.",
            "root_cause": "Historical preference for 'speed' over 'structure' during early company growth.",
            "blindness_reason": "Errors are blamed on 'human mistake' rather than 'system design'.",
            "impact": "High risk of tribal knowledge loss."
        })

    if data.approvals_per_task > 2:
        loss_points.append({
            "title": "Decision Latency Accumulation",
            "reason": f"Requires {data.approvals_per_task} distinct sign-offs, creating exponential delays.",
            "root_cause": "Lack of delegated authority; fear of making autonomous wrong decisions.",
            "blindness_reason": "Managers feel 'productive' when approving, ignoring the cost of the queue they create.",
            "impact": "Cycle time extends 4x beyond actual working hours."
        })
    
    if data.people_involved > 6:
        loss_points.append({
            "title": "Communication Overhead Entropy",
            "reason": f"{data.people_involved} active participants creates ~{data.people_involved * (data.people_involved - 1) // 2} communication pathways.",
            "root_cause": "In clear definition of 'consulted' vs 'responsible' (RACI) roles.",
            "blindness_reason": "Large meetings feel like 'collaboration', hiding the reality of 'consensus paralysis'.",
            "impact": "15% of total time lost clarifying requirements."
        })

    # Feature: Benchmarking (Mock Data per Industry Standard)
    industry_benchmark_score = random.randint(65, 85)

    # 5. Connected AI Reasoning (Gemini or Advanced Heuristics)
    import os
    import google.generativeai as genai

    # Try to use Gemini if key exists
    api_key = os.environ.get("GEMINI_API_KEY")
    ai_points = []
    
    if api_key:
        try:
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel('gemini-pro')
            
            prompt = f"""
            Act as a senior business operations consultant. Analyze this workflow description:
            "{data.description}"
            
            Context: {data.people_involved} people, {data.approvals_per_task} approvals, Tools: {', '.join(data.tools_used)}.
            
            Identify 3 specific "Invisible Loss Points". Return ONLY JSON in this format:
            [
                {{ "title": "...", "reason": "...", "root_cause": "...", "blindness_reason": "...", "impact": "..." }}
            ]
            Keep it professional, insightful, and harsh.
            """
            response = model.generate_content(prompt)
            # Basic parsing
            import json
            cleaned_response = response.text.replace('```json', '').replace('```', '')
            ai_points = json.loads(cleaned_response)
        except Exception as e:
            print(f"AI Generation failed: {e}")
    
    if ai_points:
        loss_points = ai_points
    # Else keep heuristic loss_points calculated above

    # 6. Recommendations & Scenario Simulator Data
    recommendations = []
    
    # Calculate potential savings for simulator
    savings_approval = 0
    if data.approvals_per_task > 1:
        # Simulate reducing 1 approval
        new_approvals = data.approvals_per_task - 1
        new_approval_friction = new_approvals * data.people_involved * 1.5
        new_loss = tool_friction + new_approval_friction + base_delay_impact
        savings_approval = (weekly_time_loss - new_loss) * hourly_rate
        
        recommendations.append({
            "type": "Simplify",
            "action": "Implement 'Negative Consent' protocol.",
            "impact": f"Projected Savings: ₹{int(savings_approval):,}/week",
            "simulator_action": "Reduce Approvals by 1"
        })

    if len(data.tools_used) > 3:
        recommendations.append({
            "type": "Automate",
            "action": "Unify Data Ingestion Layer.",
            "impact": "Recovers 10+ hours/week per person.",
             "simulator_action": "Consolidate Tools"
        })

    if data.people_involved > 5:
         recommendations.append({
            "type": "Restructure",
            "action": "Split into Two-Pizza Teams.",
            "impact": "Reduces coordination overhead by 40%.",
             "simulator_action": "Reduce Team Size"
        })
    
    analysis_result = {
        "weekly_time_loss_hours": round(weekly_time_loss, 1),
        "estimated_financial_loss": round(financial_loss_weekly, 2),
        "confidence_interval": {
            "lower": round(loss_lower_bound, 2),
            "upper": round(loss_upper_bound, 2)
        },
        "decision_delay_index": decision_delay_index,
        "industry_benchmark_score": industry_benchmark_score,
        "rework_loss_hours": round(rework_impact, 1),
        "waste_ratio": waste_ratio,
        "total_investment": round(total_annual_investment, 2),
        "severity": severity,
        "clarity_score": clarity,
        "invisible_loss_points": loss_points,
        "recommendations": recommendations or [{
            "type": "Review", 
            "action": "Conduct a detailed process audit.", 
            "impact": "Uncover further hidden constraints."
        }]
    }

    # Save to DB
    new_workflow = WorkflowDB(
        name=data.name,
        description=data.description,
        created_at=datetime.now().isoformat(),
        input_data=data.dict(),
        result_data=analysis_result
    )
    db.add(new_workflow)
    db.commit()
    db.refresh(new_workflow)
    
    analysis_result["id"] = new_workflow.id

    return analysis_result

@app.get("/export/pptx/{workflow_id}")
def export_pptx(workflow_id: int, db: Session = Depends(get_db)):
    workflow = db.query(WorkflowDB).filter(WorkflowDB.id == workflow_id).first()
    if not workflow:
        raise HTTPException(status_code=404, detail="Workflow not found")
        
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Invisible Business Loss Report"
    subtitle.text = f"Analysis for: {workflow.name}\nGenerated by IBLD Executive AI"

    # Slide 2: Executive Summary
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Executive Summary"
    content = slide.placeholders[1]
    
    res = workflow.result_data
    inp = workflow.input_data
    
    summary_text = (
        f"Process: {inp['name']}\n"
        f"Team Size: {inp['people_involved']} | Approvals: {inp['approvals_per_task']}\n\n"
        f"Financial Impact: ₹{res['estimated_financial_loss']:,} / week\n"
        f"Time Wasted: {res['weekly_time_loss_hours']} hours / week\n"
        f"Clarity Score: {res['clarity_score']}/100 (Industry Avg: {res['industry_benchmark_score']})"
    )
    content.text = summary_text

    # Slide 3: Invisible Loss Points
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Identified Friction Points"
    content = slide.placeholders[1]
    
    bullets = ""
    for point in res.get('invisible_loss_points', [])[:3]:
        bullets += f"• {point['title']}: {point['reason']} (Impact: {point['impact']})\n"
    content.text = bullets

    # Slide 4: Recommendations
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Strategic Action Plan"
    content = slide.placeholders[1]
    
    recs = ""
    for rec in res.get('recommendations', []):
        recs += f"• {rec['type'].upper()}: {rec['action']} -> {rec['impact']}\n"
    content.text = recs

    # Save
    filename = f"IBLD_Report_{workflow_id}.pptx"
    file_path = f"./{filename}"
    prs.save(file_path)
    
    return FileResponse(path=file_path, filename=filename, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')

from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from datetime import datetime

@app.get("/export/pdf/{workflow_id}")
def export_pdf(workflow_id: int, db: Session = Depends(get_db)):
    workflow = db.query(WorkflowDB).filter(WorkflowDB.id == workflow_id).first()
    if not workflow:
        raise HTTPException(status_code=404, detail="Workflow not found")

    res = workflow.result_data
    inp = workflow.input_data
    
    filename = f"LUMINA_Audit_Report_{workflow_id}_{datetime.now().strftime('%Y%m%d')}.pdf"
    clean_path = f"./{filename}"
    
    doc = SimpleDocTemplate(clean_path, pagesize=A4, rightMargin=40, leftMargin=40, topMargin=40, bottomMargin=40)
    story = []
    
    styles = getSampleStyleSheet()
    # Custom Styles
    style_title = ParagraphStyle('LuminaTitle', parent=styles['Heading1'], fontSize=42, textColor=colors.HexColor("#2d2626"), spaceAfter=20, alignment=1)
    style_subtitle = ParagraphStyle('LuminaSubtitle', parent=styles['Normal'], fontSize=16, textColor=colors.HexColor("#7a6e6e"), alignment=1, spaceAfter=60)
    style_h1 = ParagraphStyle('LuminaH1', parent=styles['Heading2'], fontSize=20, textColor=colors.HexColor("#2d2626"), spaceAfter=15, spaceBefore=20)
    style_metric = ParagraphStyle('LuminaMetric', parent=styles['Normal'], fontSize=18, textColor=colors.HexColor("#2d2626"), alignment=1)
    style_metric_label = ParagraphStyle('LuminaMetricLabel', parent=styles['Normal'], fontSize=10, textColor=colors.HexColor("#7a6e6e"), alignment=1)
    
    # --- PAGE 1: COVER ---
    story.append(Spacer(1, 150))
    story.append(Paragraph("LUMINA", style_title))
    story.append(Paragraph("OPERATIONAL INTELLIGENCE AUDIT", ParagraphStyle('Sub', parent=style_subtitle, fontSize=14, spaceAfter=10, tracking=2)))
    story.append(Spacer(1, 20))
    story.append(Paragraph(f"Subject: {inp['name']}", style_subtitle))
    story.append(Spacer(1, 150))
    story.append(Paragraph(f"Date: {datetime.now().strftime('%B %d, %Y')}", ParagraphStyle('Date', parent=styles['Normal'], alignment=1, textColor=colors.gray)))
    story.append(Paragraph("CONFIDENTIAL // INTERNAL USE ONLY", ParagraphStyle('Footer', parent=styles['Normal'], fontSize=8, alignment=1, textColor=colors.lightgrey)))
    story.append(PageBreak())
    
    # --- PAGE 2: EXECUTIVE SUMMARY ---
    story.append(Paragraph("Executive Summary", style_h1))
    story.append(Paragraph("At A Glance", styles['Heading3']))
    
    # Metrics Table
    financial_str = f"INR {res['estimated_financial_loss']:,.0f}"
    clarity_str = f"{res['clarity_score']}/100"
    efficiency_str = "OPTIMIZED" if res['clarity_score'] > 80 else ("SUB-OPTIMAL" if res['clarity_score'] > 50 else "CRITICAL")
    eff_color = colors.HexColor("#3fb950") if res['clarity_score'] > 80 else (colors.HexColor("#d29922") if res['clarity_score'] > 50 else colors.HexColor("#d94242"))
    
    data = [
        [Paragraph(financial_str, style_metric), Paragraph(clarity_str, style_metric), Paragraph(efficiency_str, ParagraphStyle('Eff', parent=style_metric, textColor=eff_color))],
        [Paragraph("Annual Capital Bleed", style_metric_label), Paragraph("Clarity Score", style_metric_label), Paragraph("Efficiency Rating", style_metric_label)]
    ]
    
    t = Table(data, colWidths=[2.3*inch, 2.3*inch, 2.3*inch])
    t.setStyle(TableStyle([
        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor("#fdf6f5")),
        ('BOX', (0,0), (-1,-1), 1, colors.white),
        ('PADDING', (0,0), (-1,-1), 20),
    ]))
    story.append(t)
    story.append(Spacer(1, 30))
    
    # Verdict
    story.append(Paragraph("The Operations Verdict", styles['Heading3']))
    verdict_text = f"This workflow is currently operating at a <b>{efficiency_str}</b> level. The analysis identified {len(res.get('invisible_loss_points', []))} major friction points contributing to a {res.get('waste_ratio',0)}% waste of total investment. Immediate remediation is recommended to recover the reported capital bleed."
    story.append(Paragraph(verdict_text, styles['Normal']))
    story.append(PageBreak())
    
    # --- PAGE 3: DEEP DIVE ---
    story.append(Paragraph("Friction Point Analysis", style_h1))
    
    table_data = [['Blindspot Issue', 'Root Cause', 'Impact']]
    for p in res.get('invisible_loss_points', []):
        table_data.append([
            Paragraph(p['title'], styles['Normal']),
            Paragraph(p['root_cause'] or p['blindness_reason'] or "Process Ambiguity", styles['BodyText']),
            Paragraph(p['impact'] or "High Latency", styles['BodyText'])
        ])
    
    t2 = Table(table_data, colWidths=[2*inch, 3*inch, 2*inch])
    t2.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#2d2626")),
        ('TEXTCOLOR', (0,0), (-1,0), colors.white),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
        ('FONTSIZE', (0,0), (-1,0), 10),
        ('BOTTOMPADDING', (0,0), (-1,0), 12),
        ('GRID', (0,0), (-1,-1), 1, colors.HexColor("#e0e0e0")),
        ('PADDING', (0,0), (-1,-1), 10),
    ]))
    story.append(t2)
    story.append(PageBreak())

    # --- PAGE 4: STRATEGIC ROADMAP ---
    story.append(Paragraph("Strategic Remediation Roadmap", style_h1))
    
    # Split Recs
    quick_wins = [r for r in res.get('recommendations', []) if r['type'] in ['Automate', 'Eliminate']]
    structural = [r for r in res.get('recommendations', []) if r['type'] not in ['Automate', 'Eliminate']]
    
    story.append(Paragraph("Phase 1: Quick Wins (0-30 Days)", styles['Heading3']))
    for qw in quick_wins:
        story.append(Paragraph(f"• <b>{qw['action']}</b>", styles['Normal']))
        story.append(Paragraph(f"  <i>Impact: {qw['impact']}</i>", ParagraphStyle('Indent', parent=styles['BodyText'], leftIndent=20)))
        story.append(Spacer(1, 10))
        
    story.append(Spacer(1, 20))
    story.append(Paragraph("Phase 2: Structural Transformation", styles['Heading3']))
    for st in structural:
        story.append(Paragraph(f"• <b>{st['action']}</b>", styles['Normal']))
        story.append(Paragraph(f"  <i>Impact: {st['impact']}</i>", ParagraphStyle('Indent', parent=styles['BodyText'], leftIndent=20)))
        story.append(Spacer(1, 10))
        
    story.append(Spacer(1, 40))
    story.append(Paragraph(f"Projected Annual Value Unlocked: ₹{res['estimated_financial_loss']*50:,.0f}", ParagraphStyle('FooterBold', parent=styles['Heading3'], textColor=colors.HexColor("#3fb950"), alignment=2)))

    doc.build(story)
    
    return FileResponse(path=clean_path, filename=filename, media_type='application/pdf')

@app.get("/export/pptx/{analysis_id}")
def export_pptx(analysis_id: int, db: Session = Depends(get_db)):
    result = db.query(WorkflowDB).filter(WorkflowDB.id == analysis_id).first()
    if not result:
        raise HTTPException(status_code=404, detail="Analysis not found")
        
    input_data = result.input_data
    res = result.result_data
    
    # Initialize PPTX
    prs = Presentation()
    
    # SLIDE 1: Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "LUMINA Intelligence Report"
    subtitle.text = f"Analysis: {input_data.get('name', 'Workflow')}\nGenerated: {datetime.now().strftime('%Y-%m-%d')}"
    
    # SLIDE 2: Executive Summary
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Executive Summary"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = f"Clarity Score: {res.get('clarity_score', 0)}/100"
    
    p = tf.add_paragraph()
    p.text = f"Weekly Time Loss: {res.get('weekly_time_loss_hours', 0)} hours"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"Estimated Annual Financial Loss: INR {res.get('estimated_financial_loss', 0) * 50:,.0f}"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"Risk Severity: {res.get('severity', 'Unknown')}"
    p.level = 1

    # SLIDE 3: Recommendations
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Strategic Roadmap"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Recommended Actions:"
    
    for rec in res.get('recommendations', [])[:5]: # Top 5
        p = tf.add_paragraph()
        p.text = f"{rec.get('type')}: {rec.get('action')}"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = f"Impact: {rec.get('impact')}"
        p.level = 2
        
    # Save
    filename = f"Lumina_Report_{analysis_id}.pptx"
    clean_path = os.path.join(os.getcwd(), filename)
    prs.save(clean_path)
    
    return FileResponse(path=clean_path, filename=filename, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')

@app.delete("/workflows/{workflow_id}")
def delete_workflow(workflow_id: int, db: Session = Depends(get_db)):
    # In a real app, verify `current_user` owns this
    wf = db.query(WorkflowDB).filter(WorkflowDB.id == workflow_id).first()
    if not wf:
        raise HTTPException(status_code=404, detail="Workflow not found")
    
    db.delete(wf)
    db.commit()
    return {"detail": "Deleted successfully"}

# --- AUTH ENDPOINTS ---
@app.post("/token", response_model=dict)
def login_for_access_token(form_data: OAuth2PasswordRequestForm = Depends(), db: Session = Depends(get_db)):
    user = db.query(UserDB).filter(UserDB.username == form_data.username).first()
    if not user or not verify_password(form_data.password, user.hashed_password):
        raise HTTPException(status_code=401, detail="Incorrect username or password", headers={"WWW-Authenticate": "Bearer"})
    access_token = create_access_token(data={"sub": user.username})
    return {"access_token": access_token, "token_type": "bearer", "role": user.role}

class UserCreate(BaseModel):
    username: str
    password: str

@app.post("/users/", response_model=dict)
def create_user(user: UserCreate, db: Session = Depends(get_db)):
    # Check if admin exists (poor man's seed)
    if not db.query(UserDB).first():
       role = "admin"
    else:
       role = "consultant"
       
    db_user = UserDB(username=user.username, hashed_password=get_password_hash(user.password), role=role)
    db.add(db_user)
    db.commit()
    db.refresh(db_user)
    return {"username": db_user.username, "role": db_user.role}

@app.get("/users/me")
def read_users_me(current_user: UserDB = Depends(get_current_user)):
    return {"username": current_user.username, "role": current_user.role}

@app.get("/")
def read_root():
    return {"status": "IBLD Backend Running"}
